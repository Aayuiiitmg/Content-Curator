from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from backend.app.schemas.slide import SlideContent
from backend.app.utils.file_utils import generate_temp_path

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def add_text(
    slide, text, x, y, w, h, size=20, bold=False, color=RGBColor(0, 0, 0), align="left"
):

    # Fix LLM generated array values
    if isinstance(size, list):
        size = size[0] if size else 20

    box = slide.shapes.add_textbox(x, y, w, h)

    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]

    if align == "center":
        p.alignment = PP_ALIGN.CENTER

    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT

    else:
        p.alignment = PP_ALIGN.LEFT

    run = p.add_run()

    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

    return box


def set_background(slide, style):

    fill = slide.background.fill
    fill.solid()

    if style == "dark":
        fill.fore_color.rgb = RGBColor(30, 39, 97)

    else:
        fill.fore_color.rgb = RGBColor(255, 255, 255)


def build_pptx(
    slides: List[SlideContent], presentation_title: str, theme_name="default"
):

    prs = Presentation()

    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    for s in slides:

        slide = prs.slides.add_slide(blank)

        set_background(slide, s.background_style)

        title_color = (
            RGBColor(255, 255, 255)
            if s.background_style == "dark"
            else RGBColor(30, 39, 97)
        )

        body_color = (
            RGBColor(255, 255, 255)
            if s.background_style == "dark"
            else RGBColor(40, 40, 40)
        )

        # title

        add_text(
            slide,
            s.title,
            Inches(0.5),
            Inches(0.3),
            Inches(12),
            Inches(0.8),
            s.title_size,
            True,
            title_color,
            s.alignment,
        )

        #
        # Layout decided by LLM
        #

        if s.layout == "title_only":

            continue

        elif s.layout == "big_stat":

            add_text(
                slide,
                s.bullets[0] if s.bullets else "",
                Inches(1),
                Inches(2),
                Inches(11),
                Inches(1),
                48,
                True,
                title_color,
                "center",
            )

        elif s.layout == "two_column":

            half = Inches(6)

            mid = len(s.bullets) // 2

            left = "\n".join("• " + x for x in s.bullets[:mid])

            right = "\n".join("• " + x for x in s.bullets[mid:])

            add_text(
                slide,
                left,
                Inches(0.7),
                Inches(1.5),
                half,
                Inches(5),
                s.body_size,
                False,
                body_color,
            )

            add_text(
                slide,
                right,
                Inches(7),
                Inches(1.5),
                half,
                Inches(5),
                s.body_size,
                False,
                body_color,
            )

        else:

            bullets = "\n\n".join("• " + b for b in s.bullets)

            add_text(
                slide,
                bullets,
                Inches(0.8),
                Inches(1.5),
                Inches(11.5),
                Inches(4),
                s.body_size,
                False,
                body_color,
            )

        # visual placeholder

        if s.visual_suggestion:

            add_text(
                slide,
                "Visual:\n" + s.visual_suggestion,
                Inches(9),
                Inches(5.8),
                Inches(3),
                Inches(0.7),
                10,
                False,
                body_color,
            )

        # notes

        if s.speaker_notes:

            slide.notes_slide.notes_text_frame.text = s.speaker_notes

    path = generate_temp_path(".pptx")

    prs.save(path)

    return path
