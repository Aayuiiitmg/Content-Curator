from pathlib import Path
from typing import List, Dict, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.schemas.slide import SlideContent
from app.utils.file_utils import generate_temp_path

# ── Color Themes ──────────────────────────────────────────────────────────────
THEMES: Dict[str, Dict] = {
    "midnight_executive": {
        "title_bg": RGBColor(0x1E, 0x27, 0x61),
        "content_bg": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0xCA, 0xDC, 0xFC),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "body_text": RGBColor(0x1E, 0x27, 0x61),
        "bullet_accent": RGBColor(0x1E, 0x27, 0x61),
    },
    "forest_moss": {
        "title_bg": RGBColor(0x2C, 0x5F, 0x2D),
        "content_bg": RGBColor(0xF5, 0xF5, 0xF5),
        "accent": RGBColor(0x97, 0xBC, 0x62),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "body_text": RGBColor(0x2C, 0x5F, 0x2D),
        "bullet_accent": RGBColor(0x97, 0xBC, 0x62),
    },
    "coral_energy": {
        "title_bg": RGBColor(0xF9, 0x61, 0x67),
        "content_bg": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x2F, 0x3C, 0x7E),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "body_text": RGBColor(0x2F, 0x3C, 0x7E),
        "bullet_accent": RGBColor(0xF9, 0x61, 0x67),
    },
    "charcoal_minimal": {
        "title_bg": RGBColor(0x36, 0x45, 0x4F),
        "content_bg": RGBColor(0xF2, 0xF2, 0xF2),
        "accent": RGBColor(0x21, 0x29, 0x21),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "body_text": RGBColor(0x36, 0x45, 0x4F),
        "bullet_accent": RGBColor(0x36, 0x45, 0x4F),
    },
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size, bold, color, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def build_pptx(
    slides: List[SlideContent],
    presentation_title: str,
    theme_name: str = "midnight_executive",
) -> Path:
    theme = THEMES.get(theme_name, THEMES["midnight_executive"])
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # Blank

    for slide_content in slides:
        slide = prs.slides.add_slide(blank_layout)
        is_title_slide = slide_content.slide_number == 1

        # Background
        bg_color = theme["title_bg"] if is_title_slide else theme["content_bg"]
        _set_bg(slide, bg_color)

        if is_title_slide:
            _build_title_slide(slide, slide_content, presentation_title, theme)
        elif slide_content.layout == "big_stat" and slide_content.bullets:
            _build_stat_slide(slide, slide_content, theme)
        else:
            _build_bullets_slide(slide, slide_content, theme)

    out_path = generate_temp_path(suffix=".pptx")
    prs.save(str(out_path))
    return out_path


def _build_title_slide(slide, content: SlideContent, presentation_title: str, theme: dict):
    colors = theme
    # Main title
    _add_textbox(
        slide,
        Inches(1), Inches(2), Inches(11.33), Inches(1.8),
        presentation_title,
        font_size=44, bold=True,
        color=colors["title_text"],
        align=PP_ALIGN.CENTER,
    )
    # Subtitle
    if content.bullets:
        _add_textbox(
            slide,
            Inches(1.5), Inches(4), Inches(10.33), Inches(1),
            content.bullets[0],
            font_size=20, bold=False, italic=True,
            color=colors["accent"],
            align=PP_ALIGN.CENTER,
        )


def _build_bullets_slide(slide, content: SlideContent, theme: dict):
    colors = theme
    # Title bar background strip
    strip = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, Inches(1.25),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = colors["title_bg"]
    strip.line.fill.background()

    # Title text
    _add_textbox(
        slide,
        Inches(0.4), Inches(0.15), Inches(12.5), Inches(1.0),
        content.title,
        font_size=28, bold=True,
        color=colors["title_text"],
    )

    # Slide number badge
    _add_textbox(
        slide,
        Inches(12.3), Inches(0.3), Inches(0.8), Inches(0.6),
        str(content.slide_number),
        font_size=11, bold=False,
        color=colors["accent"],
        align=PP_ALIGN.RIGHT,
    )

    # Bullets
    bullet_top = Inches(1.5)
    bullet_height = Inches(0.55)
    for i, bullet in enumerate(content.bullets[:6]):
        y = bullet_top + (i * bullet_height)
        # Accent dot
        dot = slide.shapes.add_shape(1, Inches(0.4), y + Inches(0.18), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors["bullet_accent"]
        dot.line.fill.background()

        _add_textbox(
            slide,
            Inches(0.65), y, Inches(12.0), bullet_height,
            bullet,
            font_size=16, bold=False,
            color=colors["body_text"],
        )

    # Speaker notes
    if content.speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = content.speaker_notes


def _build_stat_slide(slide, content: SlideContent, theme: dict):
    colors = theme
    # Title strip
    strip = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.25))
    strip.fill.solid()
    strip.fill.fore_color.rgb = colors["title_bg"]
    strip.line.fill.background()

    _add_textbox(
        slide,
        Inches(0.4), Inches(0.15), Inches(12.5), Inches(1.0),
        content.title,
        font_size=28, bold=True,
        color=colors["title_text"],
    )

    # Big stat
    stat = content.bullets[0] if content.bullets else ""
    _add_textbox(
        slide,
        Inches(1), Inches(1.5), Inches(11.33), Inches(2.2),
        stat,
        font_size=72, bold=True,
        color=colors["bullet_accent"],
        align=PP_ALIGN.CENTER,
    )

    # Supporting points
    for i, bullet in enumerate(content.bullets[1:5]):
        _add_textbox(
            slide,
            Inches(0.65), Inches(3.9) + (i * Inches(0.6)), Inches(12.0), Inches(0.55),
            f"• {bullet}",
            font_size=15, bold=False,
            color=colors["body_text"],
        )

    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes
