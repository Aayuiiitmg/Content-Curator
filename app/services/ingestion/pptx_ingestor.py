from pathlib import Path
from pptx import Presentation


def extract_text_from_pptx(file_path: Path) -> str:
    """Extract text from all slides in a PPTX file."""
    prs = Presentation(str(file_path))
    slides_text = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(parts))
    return "\n\n".join(slides_text)
